from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def load_lyrics(file_path):
    """Load lyrics from a text file, adding separators for sections."""
    with open(file_path, 'r') as file:
        lyrics = file.readlines()
    formatted_lyrics = []
    current_section = []
    
    for line in lyrics:
        line = line.strip()
        sections = set(
        ["Chorus","Verse","Pre-Chorus","Interlude","Bridge", "Ending"]+
        [f"Chorus {i}" for i in range(1, 70)] +
        [f"Pre-Chorus {i}" for i in range(1,70)] +
        [f"Bridge {i}" for i in range(1, 70)] +
        [f"Verse {i}" for i in range(1, 70)] +
        [f"Blank {i}" for i in range(1, 70)] +
        ["[Chorus]","[Verse]","[Pre-Chorus]","[Interlude]","[Bridge]", "[Turn]", "[Outro]"]+
        [f"[Chorus {i}]" for i in range(1, 70)] +
        [f"[Pre-Chorus {i}]" for i in range(1,70)] +
        [f"[Bridge {i}]" for i in range(1, 70)] +
        [f"[Verse {i}]" for i in range(1, 70)] +
        [f"[Blank {i}]" for i in range(1, 70)]
)
        if line in sections:  # Add more section titles as needed
            if current_section:  # If there's an existing section, append it
                formatted_lyrics.append(current_section)
                current_section = []
            current_section.append(line)  # Add the section title to the current section
        elif line:
            current_section.append(line)  # Add the regular line to the current section
    
    if current_section:  # Append the last section if exists
        formatted_lyrics.append(current_section)
    
    return formatted_lyrics


def create_powerpoint(lyrics, output_file):
    """Create a PowerPoint presentation with lyrics."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create a title slide with only title layout
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout (with title only)
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set the title text to the filename (without extension)
    title_text = os.path.splitext(os.path.basename(output_file))[0]  # Get the filename without extension
    title = title_slide.shapes.title
    title.text = title_text

    # Set background color of title slide to black
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color

    # Set the title text color to white
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color for title

    title.text_frame.paragraphs[0].font.name = 'Calibri Light (Headings)'  # Set font family to Calibri Light (Headings)
    title.text_frame.paragraphs[0].font.size = Pt(60)  # Set font size to 60 points

    # Center the title text horizontally and vertically
    title.left = Inches(1)  # Position from the left
    title.width = prs.slide_width - Inches(2)  # Set width to slide width minus padding
    title.top = int((prs.slide_height - title.height) / 2)  # Center vertically

    # Create lyrics slides
def create_powerpoint(lyrics_sections, output_file):
    """Create a PowerPoint presentation with lyrics."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create a title slide with only title layout
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout (with title only)
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Set the title text to the filename (without extension)
    title_text = os.path.splitext(os.path.basename(output_file))[0]  # Get the filename without extension
    title = title_slide.shapes.title
    title.text = title_text

    # Set background color of title slide to black
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color

    # Set the title text color to white
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color for title
    title.text_frame.paragraphs[0].font.name = 'Calibri Light (Headings)'  # Set font family to Calibri Light (Headings)
    title.text_frame.paragraphs[0].font.size = Pt(60)  # Set font size to 60 points

    # Center the title text horizontally and vertically
    title.left = Inches(1)  # Position from the left
    title.width = prs.slide_width - Inches(2)  # Set width to slide width minus padding
    title.top = int((prs.slide_height - title.height) / 2)  # Center vertically

    # Create lyrics slides
    for section in lyrics_sections:
        # Skip the first item in each section (the title) and process the remaining lines
        k = 3
        decrease = False
        two = False
        if len(section) == 1:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color
        for i in range(1, len(section), 3):  # Start from index 1 to skip the title

            if i == len(section) - 1:
                if decrease:
                    i -= 1
                
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            
            # Set the background color to black
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color

            if len(section) == 3:
                two = True
                k = 2

            if len(section) == 5:
                    two = True
                    decrease = True
                    k = 2
            
            # Create a textbox for lyrics
            if two == False:
                textbox = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(13), Inches(3.67))  # Adjust size and position
                text_frame = textbox.text_frame
            else:
                textbox = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(13), Inches(2.67))  # Adjust size and position
                text_frame = textbox.text_frame


            # Add up to 3 lines to the slide
            for j in range(k):
                if i + j < len(section): 
                    p = text_frame.add_paragraph()  # Create a new paragraph
                    p.text = section[i+j]  # Add the line without leading/trailing whitespace
                    p.font.size = Pt(50)  # Adjust font size as needed
                    p.font.name = "Calibri Light"  # Use Calibri Light
                    p.font.color.rgb = RGBColor(255, 255, 255)  # Set text color to white
                    p.alignment = PP_ALIGN.CENTER  # Center align the text

                # Center the lyrics text box
                textbox.left = Inches(0)  # Position from the left
                textbox.width = prs.slide_width  # Set width to slide width
                textbox.top = int((prs.slide_height - textbox.height) / 2)  # Center vertically

                if i + j + 5 == len(section):
                    if len(section) % 3 == 2:
                        two = True
                        decrease = True
                        k = 2
                else:
                    k = 3

                if i + j + 4 == len(section):
                    if len(section) % 3 == 0:
                        two = True

                



    prs.save(output_file)
    print(f"PowerPoint saved as {output_file}")

# Usage
lyrics_file_path = 'Large_Group_10-18-2024.txt'  # Path to your lyrics text file
output_file = 'Large_Group_10-18-2024.pptx'  # Output PowerPoint file

lyrics_sections = load_lyrics(lyrics_file_path)  # Load the lyrics
create_powerpoint(lyrics_sections, output_file)  # Create the PowerPoint